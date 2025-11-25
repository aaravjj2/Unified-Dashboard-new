"""Build presentation PPTX from YAML slide definition.

Usage:
    pip install python-pptx pyyaml
    python presentation/build_presentation.py slides_definition.yaml SharkTank_Pitch.pptx

This script reads `slides_definition.yaml` and produces a PowerPoint with speaker notes.
"""
import sys
import yaml
from pptx import Presentation
from pptx.util import Inches, Pt
import os

SLIDE_LAYOUT_TITLE = 0
SLIDE_LAYOUT_TITLE_AND_CONTENT = 1
SLIDE_LAYOUT_BLANK = 6

BULLET_FONT_SIZE = Pt(22)
NOTE_FONT_SIZE = Pt(12)


def add_title_slide(prs, slide_def):
    slide = prs.slides.add_slide(prs.slide_layouts[SLIDE_LAYOUT_TITLE])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = slide_def.get('title', '')
    subtitle.text = slide_def.get('subtitle', '')
    # speaker notes
    notes = slide.notes_slide.notes_text_frame
    notes.text = slide_def.get('notes', '')
    # Optional image
    img = slide_def.get('image') or slide_def.get('images')
    if img:
        # allow images: single path or list
        path = img[0] if isinstance(img, list) else img
        if os.path.exists(path):
            try:
                pic = slide.shapes.add_picture(path, Inches(0.5), Inches(1.6), width=Inches(9))
                pic.shadow.inherit = False
            except Exception:
                pass


def add_bullets_slide(prs, slide_def):
    slide = prs.slides.add_slide(prs.slide_layouts[SLIDE_LAYOUT_TITLE_AND_CONTENT])
    title = slide.shapes.title
    title.text = slide_def.get('title', '')
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    bullets = slide_def.get('bullets', [])
    if bullets:
        for i, b in enumerate(bullets):
            if i == 0:
                p = body.paragraphs[0]
            else:
                p = body.add_paragraph()
            p.text = b
            p.level = 0
            p.font.size = BULLET_FONT_SIZE
    notes = slide_def.get('notes', '')
    slide.notes_slide.notes_text_frame.text = notes
    # Optional images for content slides (placed to the right)
    img = slide_def.get('image') or slide_def.get('images')
    if img:
        path = img[0] if isinstance(img, list) else img
        if os.path.exists(path):
            try:
                # place image on the right column
                left = Inches(6.5)
                top = Inches(1.6)
                width = Inches(3.0)
                slide.shapes.add_picture(path, left, top, width=width)
            except Exception:
                pass


def add_image_only_slide(prs, slide_def):
    # Blank layout with just an image (good for large visuals)
    slide = prs.slides.add_slide(prs.slide_layouts[SLIDE_LAYOUT_BLANK])
    img = slide_def.get('image') or slide_def.get('images')
    if not img:
        return
    path = img[0] if isinstance(img, list) else img
    if os.path.exists(path):
        try:
            # Fit image to slide width
            slide_w = prs.slide_width
            pic = slide.shapes.add_picture(path, Inches(0), Inches(0))
            # scale width to slide width if larger
            if pic.width > slide_w:
                scale = slide_w / pic.width
                pic.width = int(pic.width * scale)
                pic.height = int(pic.height * scale)
        except Exception:
            pass


def add_appendix_slide(prs, slide_def):
    # Reuse bullets layout
    add_bullets_slide(prs, slide_def)


def build_presentation(yaml_path, out_pptx):
    with open(yaml_path, 'r', encoding='utf-8') as f:
        doc = yaml.safe_load(f)

    slides = doc.get('slides', [])
    prs = Presentation()

    for s in slides:
        t = s.get('type')
        if t == 'title':
            add_title_slide(prs, s)
        elif t == 'bullets':
            # if slide requests a full-image visual, use image-only layout
            use_full_image = s.get('full_image', False) or (s.get('image') and s.get('image_only', False))
            if use_full_image:
                add_image_only_slide(prs, s)
                # still add notes
                prs.slides[-1].notes_slide.notes_text_frame.text = s.get('notes', '')
            else:
                add_bullets_slide(prs, s)
        elif t == 'appendix':
            add_appendix_slide(prs, s)
        else:
            # default to bullets
            add_bullets_slide(prs, s)

    prs.save(out_pptx)
    print(f"Saved presentation to {out_pptx}")


if __name__ == '__main__':
    if len(sys.argv) < 3:
        print("Usage: python build_presentation.py slides_definition.yaml out.pptx")
        sys.exit(1)
    yaml_path = sys.argv[1]
    out = sys.argv[2]
    build_presentation(yaml_path, out)
